#!/usr/bin/env python3
"""
Howatson+Company Template - Example Implementation
Demonstrates how to create slides using python-pptx following exact specifications.

CRITICAL IMPLEMENTATION NOTES:
- Font name is "VC Nudge Black" for headlines
- All shapes require explicit shadow removal
- Multi-line text requires formatting EACH paragraph
- Sidebar color ALWAYS matches heading color
- Use connector lines for underlines (not rectangles)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ============================================================================
# CONSTANTS
# ============================================================================

HEADLINE_FONT = "VC Nudge Black"
BODY_FONT = "Manrope"
BODY_FONT_LIGHT = "Manrope Light"

COLORS = {
    'black': RGBColor(0, 0, 0),
    'white': RGBColor(255, 255, 255),
    'beige': RGBColor(237, 218, 193),
    'grey': RGBColor(224, 224, 224),
    'cyan': RGBColor(118, 219, 236),
    'blue': RGBColor(61, 114, 254),
    'lime': RGBColor(224, 234, 83),
    'red': RGBColor(255, 79, 77),
    'pink': RGBColor(233, 114, 226),
    'green': RGBColor(0, 185, 72),
}

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def remove_shadow(shape):
    """Remove ALL shadow effects from any shape.

    CRITICAL: Call this on EVERY shape you create:
    - Text boxes
    - Rectangles (image placeholders)
    - Connector lines (underlines)
    - Any other shapes
    """
    spPr = shape._element.spPr

    # Remove effectLst (contains shadow definitions)
    for effect in spPr.findall(qn('a:effectLst')):
        spPr.remove(effect)

    # Remove outerShdw directly if nested
    for shadow in spPr.findall('.//' + qn('a:outerShdw')):
        shadow.getparent().remove(shadow)

    # Remove effect style references
    for style in spPr.findall(qn('a:effectRef')):
        spPr.remove(style)


def add_line(slide, x1, y1, x2, y2, color, width_pt=1):
    """Add a simple connector line with no shadow.

    Use this for underlines instead of rectangles.
    """
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    remove_shadow(connector)
    return connector


def add_headline(slide, left, top, width, height, lines, font_size, color,
                 alignment=PP_ALIGN.LEFT, line_spacing=0.85):
    """Add multi-line headline text with proper formatting.

    CRITICAL: Formats EACH paragraph individually.
    Using \\n in text creates separate paragraphs, each needs formatting.

    Args:
        lines: Either a string (single line) or list of strings (multi-line)
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = False

    if isinstance(lines, str):
        lines = [lines]

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = HEADLINE_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.alignment = alignment
        p.line_spacing = line_spacing

    remove_shadow(txBox)
    return txBox


def add_body_text(slide, left, top, width, height, text, color,
                  font_size=14, line_spacing=1.4):
    """Add body text with consistent formatting across all paragraphs.

    CRITICAL: Formats EACH paragraph, not just the first.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    # Split by double newline for paragraphs
    paragraphs = text.split('\n\n') if '\n\n' in text else [text]

    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = para_text
        p.font.name = BODY_FONT_LIGHT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.line_spacing = line_spacing
        p.space_after = Pt(font_size)

    remove_shadow(txBox)
    return txBox


def add_sidebar(slide, color=None):
    """Add vertical sidebar with branding.

    CRITICAL RULES:
    1. Line runs Y=0 to Y=7.5 (FULL slide height, edge-to-edge)
    2. Color ALWAYS matches heading color:
       - White/light backgrounds → BLACK
       - Black backgrounds → WHITE
       - Colored backgrounds (cyan, blue, lime) → BLACK (NOT white!)
    """
    if color is None:
        color = COLORS['black']

    # FULL HEIGHT: Y=0 to Y=7.5 (no gaps!)
    add_line(slide, 0.35, 0, 0.35, 7.5, color, 1)

    # Rotated branding text at bottom
    txBox = slide.shapes.add_textbox(
        Inches(0.05), Inches(6.2),
        Inches(0.25), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.text = "© HOWATSON+COMPANY"
    tf.word_wrap = False

    p = tf.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(7)
    p.font.color.rgb = color

    txBox.rotation = 270
    remove_shadow(txBox)


def set_background_color(slide, color):
    """Set solid colour background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================================================
# SLIDE CREATION FUNCTIONS
# ============================================================================

def create_cover_slide(prs, client_name="CLIENT NAME", presentation_name=None):
    """
    Create Cover Slide (Layout 1)
    Black background with large white text and central image.
    """
    if presentation_name is None:
        presentation_name = ["PRESENTATION", "NAME"]
    elif isinstance(presentation_name, str):
        presentation_name = presentation_name.split('\n')

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set black background
    set_background_color(slide, COLORS['black'])

    # Add vertical branding (white on black)
    add_sidebar(slide, COLORS['white'])

    # Client name (top left)
    add_headline(
        slide, 0.7, 2.69, 5.0, 0.6,
        f"+ {client_name}",
        28, COLORS['white']
    )

    # Presentation name (top right) - multi-line
    add_headline(
        slide, 8.5, 2.54, 4.5, 1.2,
        presentation_name,
        28, COLORS['white'],
        alignment=PP_ALIGN.RIGHT
    )

    # Note: Large "HOWATSON+COMPANY" text would be added here
    # along with a central image. See template-specs.md for exact measurements.

    return slide


def create_section_divider(prs, section_number="01", section_title=None, bg_color='cyan'):
    """
    Create Section Divider Slide (Layout 3)
    Accent colour background with large number and title.

    CRITICAL:
    - Section title uses MSO_ANCHOR.BOTTOM (not MIDDLE)
    - Sidebar is BLACK on colored backgrounds
    """
    if section_title is None:
        section_title = ["SECTION", "TITLE"]
    elif isinstance(section_title, str):
        section_title = section_title.split('\n') if '\n' in section_title else [section_title]

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background colour
    set_background_color(slide, COLORS[bg_color])

    # CRITICAL: Sidebar is BLACK on colored backgrounds (NOT white!)
    add_sidebar(slide, COLORS['black'])

    # Section number (large, bleeds left)
    add_headline(
        slide, -2.0, -0.5, 10.0, 8.5,
        section_number,
        500, COLORS['black']
    )

    # Section title (right side) - BOTTOM aligned
    txBox = slide.shapes.add_textbox(
        Inches(5.5), Inches(4.0),
        Inches(7.5), Inches(3.0)
    )
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM  # CRITICAL: Bottom-aligned
    tf.word_wrap = False

    for i, line in enumerate(section_title):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = HEADLINE_FONT
        p.font.size = Pt(85)
        p.font.color.rgb = COLORS['black']
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 0.85

    remove_shadow(txBox)

    return slide


def create_three_column_content(prs, headline="HEADLINE GOES HERE", columns_data=None):
    """
    Create 3-Column Content Slide (Layout 5)
    White background with headline and three equal columns.

    columns_data: List of dicts with 'header' and 'body' keys
    """
    if columns_data is None:
        columns_data = [
            {'header': 'Heading 1', 'body': 'Lorem ipsum dolor sit amet...'},
            {'header': 'Heading 2', 'body': 'Lorem ipsum dolor sit amet...'},
            {'header': 'Heading 3', 'body': 'Lorem ipsum dolor sit amet...'},
        ]

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # White background
    set_background_color(slide, COLORS['white'])

    # Add sidebar (black on white background)
    add_sidebar(slide, COLORS['black'])

    # Main headline - handle multi-line
    headline_lines = headline.split('\n') if '\n' in headline else [headline]
    add_headline(
        slide, 0.7, 0.7, 12.0, 2.5,
        headline_lines,
        70, COLORS['black']
    )

    # Three columns
    column_width = 3.96
    column_start_y = 3.8
    gutter = 0.5

    for i, col_data in enumerate(columns_data[:3]):
        col_x = 0.7 + i * (column_width + gutter)

        # Column header
        header_box = slide.shapes.add_textbox(
            Inches(col_x), Inches(column_start_y),
            Inches(column_width), Inches(0.4)
        )
        tf = header_box.text_frame
        tf.text = col_data['header']

        p = tf.paragraphs[0]
        p.font.name = BODY_FONT
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['black']

        remove_shadow(header_box)

        # Underline - use CONNECTOR line (not rectangle!)
        add_line(
            slide,
            col_x, column_start_y + 0.45,
            col_x + column_width, column_start_y + 0.45,
            COLORS['black'], 1
        )

        # Body text
        add_body_text(
            slide,
            col_x, column_start_y + 0.6,
            column_width, 2.5,
            col_data['body'],
            COLORS['black'],
            font_size=16
        )

    return slide


def create_split_image_text(prs, headline="HEADING 1", body_text="", image_path=None):
    """
    Create Split Image/Text Slide (Layout 9)
    Left 40% text, right 60% image.
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # White background
    set_background_color(slide, COLORS['white'])

    # Add sidebar (black on white)
    add_sidebar(slide, COLORS['black'])

    # Headline (left side) - handle multi-line
    headline_lines = headline.split('\n') if '\n' in headline else [headline]
    add_headline(
        slide, 0.7, 1.2, 5.0, 2.0,
        headline_lines,
        70, COLORS['black']
    )

    # Body text (left side)
    add_body_text(
        slide, 0.7, 3.5, 4.8, 2.5,
        body_text,
        COLORS['black'],
        font_size=17
    )

    # Image placeholder (right side)
    if image_path:
        left = Inches(5.8)
        top = Inches(0.36)
        width = Inches(7.5)
        pic = slide.shapes.add_picture(image_path, left, top, width=width)
        remove_shadow(pic)
    else:
        # Add placeholder rectangle
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.8), Inches(0.36),
            Inches(7.5), Inches(6.78)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = COLORS['beige']
        placeholder.line.fill.background()
        remove_shadow(placeholder)

    return slide


# ============================================================================
# MAIN EXAMPLE
# ============================================================================

def create_example_presentation():
    """Create a sample presentation using the Howatson template."""

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    create_cover_slide(
        prs,
        client_name="ACME CORP",
        presentation_name=["STRATEGY", "DECK"]
    )

    # Slide 2: Section Divider (cyan)
    create_section_divider(
        prs,
        section_number="01",
        section_title=["INTRODUCTION"],
        bg_color='cyan'
    )

    # Slide 3: Three-column content
    columns = [
        {
            'header': 'Challenge',
            'body': 'Lorem ipsum dolor sit amet, consectetur adipiscing elit. '
                   'Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.'
        },
        {
            'header': 'Approach',
            'body': 'Ut enim ad minim veniam, quis nostrud exercitation ullamco '
                   'laboris nisi ut aliquip ex ea commodo consequat.'
        },
        {
            'header': 'Results',
            'body': 'Duis aute irure dolor in reprehenderit in voluptate velit '
                   'esse cillum dolore eu fugiat nulla pariatur.'
        }
    ]
    create_three_column_content(
        prs,
        headline="OUR STRATEGIC FRAMEWORK",
        columns_data=columns
    )

    # Slide 4: Split image/text
    create_split_image_text(
        prs,
        headline="KEY INSIGHT",
        body_text="Moving forward we can keep things simple with small font and "
                 "lots of white space. This helps to make the deck feel fresh and clean.\n\n"
                 "Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do "
                 "eiusmod tempor incididunt ut labore et dolore magna aliqua."
    )

    # Slide 5: Another section divider (blue)
    create_section_divider(
        prs,
        section_number="02",
        section_title=["RECOMMENDATIONS"],
        bg_color='blue'
    )

    # Save presentation
    output_path = "/mnt/user-data/outputs/howatson_example.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_example_presentation()
    print("\nExample presentation created successfully!")
    print("Review the output to see how the template is implemented.")
