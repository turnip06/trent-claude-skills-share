#!/usr/bin/env python3
"""
Generate an HCO presentation from a mapping dictionary.

V2 IMPROVEMENTS:
- Uses python-pptx-text-replacer for formatting-safe text replacement
- Fixed slide duplication that preserves images and relationships
- Better error handling and validation
- Font property preservation as fallback

Usage:
    python scripts/generate_presentation.py

Functions:
    duplicate_slide(prs, index) - Duplicate a template slide with relationships
    replace_placeholder(slide, placeholder, new_text) - Replace across entire slide
    replace_text_preserve_formatting(slide, replacements) - Batch replace with formatting
"""

import sys
import os
import re
from copy import deepcopy
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Try to import the text replacer package
try:
    from python_pptx_text_replacer import TextReplacer
    HAS_TEXT_REPLACER = True
except ImportError:
    HAS_TEXT_REPLACER = False
    print("Warning: python-pptx-text-replacer not installed. Using fallback method.")
    print("Install with: pip install python-pptx-text-replacer")

# Default paths (relative to project root)
TEMPLATE_PATH = "assets/ppt-template-v2.pptx"
OUTPUT_PATH = "output.pptx"

# VC Nudge Black line spacing (0.7 = 70% of single line spacing)
VC_NUDGE_LINE_SPACING = 0.7


def generate_output_filename(client_name: str = "", presentation_title: str = "") -> str:
    """
    Generate a descriptive output filename based on content.

    Creates filesystem-safe filename from client name and presentation title.
    Falls back to "output.pptx" if no content provided.

    Args:
        client_name: Client name (e.g., "+ACME CORP" or "Acme")
        presentation_title: Presentation title (e.g., "Q4 Brand Strategy")

    Returns:
        Sanitized filename with .pptx extension

    Examples:
        generate_output_filename("+ACME CORP", "Q4 BRAND STRATEGY")
        -> "acme-corp_q4-brand-strategy.pptx"

        generate_output_filename("Acme", "")
        -> "acme.pptx"

        generate_output_filename("", "")
        -> "output.pptx"
    """
    def sanitize(text: str) -> str:
        """Convert text to filesystem-safe lowercase slug."""
        # Remove leading + (HCO brand convention)
        text = text.lstrip('+')
        # Convert to lowercase
        text = text.lower()
        # Replace spaces and underscores with hyphens
        text = re.sub(r'[\s_]+', '-', text)
        # Remove any characters that aren't alphanumeric or hyphens
        text = re.sub(r'[^a-z0-9-]', '', text)
        # Collapse multiple hyphens
        text = re.sub(r'-+', '-', text)
        # Strip leading/trailing hyphens
        text = text.strip('-')
        return text

    client_slug = sanitize(client_name)
    title_slug = sanitize(presentation_title)

    if client_slug and title_slug:
        return f"{client_slug}_{title_slug}.pptx"
    elif client_slug:
        return f"{client_slug}.pptx"
    elif title_slug:
        return f"{title_slug}.pptx"
    else:
        return OUTPUT_PATH  # Fallback to default


def get_font_properties(run: Any) -> dict[str, Any]:
    """
    Extract font properties from a run for preservation.

    Args:
        run: A pptx run object

    Returns:
        dict with font properties
    """
    font = run.font
    props = {
        'name': font.name,
        'size': font.size,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
    }

    # Handle color carefully - it may not be set
    try:
        if font.color and font.color.type is not None:
            props['color_rgb'] = font.color.rgb
        else:
            props['color_rgb'] = None
    except Exception:
        props['color_rgb'] = None

    return props


def apply_font_properties(run: Any, props: dict[str, Any]) -> None:
    """
    Apply preserved font properties to a run.

    Args:
        run: A pptx run object
        props: dict of font properties from get_font_properties()
    """
    font = run.font

    if props.get('name'):
        font.name = props['name']
    if props.get('size'):
        font.size = props['size']
    if props.get('bold') is not None:
        font.bold = props['bold']
    if props.get('italic') is not None:
        font.italic = props['italic']
    if props.get('underline') is not None:
        font.underline = props['underline']
    if props.get('color_rgb'):
        font.color.rgb = props['color_rgb']


def get_paragraph_properties(para: Any) -> dict[str, Any]:
    """
    Extract paragraph properties for preservation.

    Args:
        para: A pptx paragraph object

    Returns:
        dict with paragraph properties including line_spacing
    """
    return {
        'line_spacing': para.line_spacing,
        'space_before': para.space_before,
        'space_after': para.space_after,
    }


def apply_paragraph_properties(para: Any, props: dict[str, Any]) -> None:
    """
    Apply preserved paragraph properties.

    Args:
        para: A pptx paragraph object
        props: dict of paragraph properties from get_paragraph_properties()
    """
    if props.get('line_spacing') is not None:
        para.line_spacing = props['line_spacing']
    if props.get('space_before') is not None:
        para.space_before = props['space_before']
    if props.get('space_after') is not None:
        para.space_after = props['space_after']


def is_vc_nudge_font(para: Any) -> bool:
    """
    Check if paragraph uses VC Nudge font.

    Args:
        para: A pptx paragraph object

    Returns:
        True if any run in the paragraph uses VC Nudge font
    """
    for run in para.runs:
        font_name = run.font.name
        if font_name and 'VC Nudge' in font_name:
            return True
    return False


def ensure_vc_nudge_line_spacing(para: Any) -> None:
    """
    Ensure VC Nudge paragraphs have 0.7 line spacing.

    This is called after text replacement to ensure the tight line spacing
    is preserved even if it was lost during replacement.

    Args:
        para: A pptx paragraph object
    """
    if is_vc_nudge_font(para):
        para.line_spacing = VC_NUDGE_LINE_SPACING


def adjust_heading_height(shape, new_text, max_chars_per_line=25):
    """
    Expand text box height for long headlines to prevent clipping.

    When text exceeds the original box height, this expands the box upward
    (since headings typically use bottom anchoring) to accommodate more lines.

    Args:
        shape: Shape object with text frame
        new_text: The text being inserted
        max_chars_per_line: Estimated characters per line for this font/width

    Returns:
        True if height was adjusted, False otherwise
    """
    if not shape.has_text_frame:
        return False

    # Calculate approximate lines needed
    text_length = len(new_text)
    estimated_lines = max(1, (text_length + max_chars_per_line - 1) // max_chars_per_line)

    # Minimum height per line (approximately 0.5in for large headlines)
    min_height_per_line = Inches(0.5)
    needed_height = min_height_per_line * estimated_lines

    # Only expand if needed height exceeds current height
    if needed_height > shape.height:
        # Calculate how much to expand
        height_increase = needed_height - shape.height

        # Move shape up to expand upward (preserve bottom position)
        shape.top = shape.top - height_increase
        shape.height = needed_height
        return True

    return False


def adjust_section_title_position(slide, section_number):
    """
    Adjust section title position to prevent overlap with wide section numbers.

    Numbers like "02", "03", "08" use wider glyphs than "01", "11", etc.
    This shifts the title text box right when needed.

    Args:
        slide: Slide object
        section_number: The section number string (e.g., "01", "02")

    Returns:
        True if position was adjusted, False otherwise
    """
    # Glyphs that appear wider at large font sizes and require title position adjustment
    # 0,2,3,6,8,9 are rounder/wider than 1,4,5,7 in VC Nudge Black at 120pt
    WIDE_SECTION_NUMBER_GLYPHS = frozenset('023689')

    # Calculate width adjustment based on number of wide glyphs
    num_str = str(section_number).zfill(2)
    wide_count = sum(1 for c in num_str if c in WIDE_SECTION_NUMBER_GLYPHS)

    if wide_count == 0:
        return False  # No adjustment needed

    # Find the title shape (contains section.title placeholder or the actual title)
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            # Check if this is the title shape (either has placeholder or is positioned correctly)
            if '{{section.title}}' in text or (shape.left > Inches(4) and shape.top < Inches(2)):
                # Add offset: approximately 0.25in per wide glyph
                adjustment = Inches(0.25 * wide_count)
                shape.left = shape.left + adjustment
                return True

    return False


def duplicate_slide(prs: Presentation, index: int) -> Any:
    """
    Duplicate a slide from the template, preserving backgrounds, images, and all relationships.

    V3 IMPROVEMENTS:
    - Copies entire cSld element (background + shapes) instead of just shapes
    - Builds rId mapping to fix broken image references
    - Updates all rId attributes in copied elements to use new relationship IDs

    Args:
        prs: Presentation object
        index: 0-based index of slide to duplicate

    Returns:
        The newly created slide
    """
    if index >= len(prs.slides):
        raise IndexError(f"Slide index {index} exceeds presentation size ({len(prs.slides)} slides)")

    template = prs.slides[index]
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Build rId mapping: old_rId -> new_rId
    # This is critical for fixing broken image references
    rid_map = {}
    for rel_key, rel in template.part.rels.items():
        # Skip relationship types that shouldn't be copied
        if rel.reltype == RT.NOTES_SLIDE:
            continue
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue

        # For internal relationships (images, charts, etc.)
        if not rel.is_external:
            try:
                target = rel.target_part
                new_rid = new_slide.part.relate_to(target, rel.reltype)
                rid_map[rel_key] = new_rid
            except Exception as e:
                print(f"Warning: Could not copy internal relationship {rel_key}: {e}")
        else:
            # For external relationships (hyperlinks, etc.)
            try:
                new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                rid_map[rel_key] = new_rid
            except Exception as e:
                print(f"Warning: Could not copy external relationship {rel_key}: {e}")

    # Copy entire cSld element (includes background AND shapes)
    # This fixes the cyan background issue
    cSld_copy = deepcopy(template.element.cSld)

    # Update all rId references in the copied element to use new relationship IDs
    # XML namespaces for relationship attributes
    nsmap = {
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }

    # Find and update all elements with relationship ID attributes
    for elem in cSld_copy.iter():
        # Check for r:embed (embedded images, charts)
        embed_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
        if embed_attr in elem.attrib:
            old_rid = elem.attrib[embed_attr]
            if old_rid in rid_map:
                elem.attrib[embed_attr] = rid_map[old_rid]

        # Check for r:link (linked media)
        link_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link'
        if link_attr in elem.attrib:
            old_rid = elem.attrib[link_attr]
            if old_rid in rid_map:
                elem.attrib[link_attr] = rid_map[old_rid]

        # Check for r:id (hyperlinks and other relationships)
        id_attr = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        if id_attr in elem.attrib:
            old_rid = elem.attrib[id_attr]
            if old_rid in rid_map:
                elem.attrib[id_attr] = rid_map[old_rid]

    # Replace new slide's cSld with our updated copy
    # First remove the existing cSld
    existing_cSld = new_slide.element.cSld
    new_slide.element.remove(existing_cSld)

    # Insert the copied cSld at the correct position (first child)
    new_slide.element.insert(0, cSld_copy)

    return new_slide


def replace_in_shape_with_formatting(shape: Any, placeholder: str, new_text: str) -> bool:
    """
    Replace placeholder text in a shape while preserving formatting.

    This handles:
    - Run fragmentation (placeholder split across runs)
    - Font property preservation
    - Paragraph property preservation (including line_spacing)
    - VC Nudge 0.7 line spacing enforcement
    - Grouped shapes (recursive)

    Args:
        shape: Shape object
        placeholder: Placeholder token (e.g., "{{title.client_name}}")
        new_text: Text to replace with

    Returns:
        True if replacement was made, False otherwise
    """
    replaced = False

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            # First, check if placeholder exists in full paragraph text
            full_text = ''.join(run.text for run in para.runs)

            if placeholder in full_text:
                # Placeholder found - now handle the replacement carefully

                # Capture paragraph properties before replacement
                para_props = get_paragraph_properties(para)

                # Simple case: placeholder is entirely within one run
                for run in para.runs:
                    if placeholder in run.text:
                        # Preserve formatting
                        props = get_font_properties(run)
                        run.text = run.text.replace(placeholder, new_text)
                        apply_font_properties(run, props)
                        replaced = True

                # Complex case: placeholder spans multiple runs
                if not replaced and placeholder in full_text:
                    # Find which runs contain parts of the placeholder
                    placeholder_start = full_text.find(placeholder)
                    placeholder_end = placeholder_start + len(placeholder)

                    # Track position as we iterate through runs
                    current_pos = 0
                    runs_to_modify = []

                    for i, run in enumerate(para.runs):
                        run_start = current_pos
                        run_end = current_pos + len(run.text)

                        # Check if this run overlaps with placeholder
                        if run_end > placeholder_start and run_start < placeholder_end:
                            runs_to_modify.append({
                                'index': i,
                                'run': run,
                                'start': run_start,
                                'end': run_end
                            })

                        current_pos = run_end

                    if runs_to_modify:
                        # Get formatting from first affected run
                        first_run = runs_to_modify[0]['run']
                        props = get_font_properties(first_run)

                        # Clear all affected runs except first
                        for run_info in runs_to_modify[1:]:
                            run_info['run'].text = ''

                        # Replace in first run (reconstruct the text)
                        before_placeholder = full_text[:placeholder_start]
                        after_placeholder = full_text[placeholder_end:]

                        # Put everything in first run
                        first_run_start = runs_to_modify[0]['start']
                        first_run.text = full_text[first_run_start:placeholder_start] + new_text

                        # If there's text after, put it in the last run
                        if after_placeholder and len(runs_to_modify) > 1:
                            last_run = runs_to_modify[-1]['run']
                            last_run.text = after_placeholder[:(runs_to_modify[-1]['end'] - placeholder_end)]

                        apply_font_properties(first_run, props)
                        replaced = True

                # Restore paragraph properties after replacement
                if replaced:
                    apply_paragraph_properties(para, para_props)
                    # Ensure VC Nudge headlines have 0.7 line spacing
                    ensure_vc_nudge_line_spacing(para)

    # Handle grouped shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            if replace_in_shape_with_formatting(child, placeholder, new_text):
                replaced = True

    return replaced


def replace_placeholder(slide: Any, placeholder: str, new_text: str) -> bool:
    """
    Replace a placeholder across all shapes in a slide.

    Uses formatting-preserving replacement method.

    Args:
        slide: Slide object
        placeholder: Placeholder token (e.g., "{{title.client_name}}")
        new_text: Text to replace with

    Returns:
        True if replacement was made, False otherwise
    """
    replaced = False
    for shape in slide.shapes:
        if replace_in_shape_with_formatting(shape, placeholder, new_text):
            replaced = True
    return replaced


def replace_text_batch(prs_path, output_path, replacements):
    """
    Batch replace text using python-pptx-text-replacer for best formatting preservation.

    This is the recommended method when making multiple replacements.

    Args:
        prs_path: Path to input presentation
        output_path: Path for output presentation
        replacements: Dict mapping placeholder -> new_text

    Returns:
        Path to output file
    """
    if not HAS_TEXT_REPLACER:
        raise ImportError("python-pptx-text-replacer required for batch replacement")

    replacer = TextReplacer(
        prs_path,
        tables=True,
        charts=True,
        textframes=True
    )

    # Convert dict to list of tuples for the replacer
    replacement_list = [(k, v) for k, v in replacements.items()]

    replacer.replace_text(replacement_list)
    replacer.write_presentation_to_file(output_path)

    return output_path


def generate_presentation(content, template_path=TEMPLATE_PATH, output_path=OUTPUT_PATH,
                          use_batch_replacer=True):
    """
    Generate a presentation from content dictionary.

    Args:
        content: Dict mapping slide_index -> {placeholder: text} pairs
                 Example: {0: {"{{title.client_name}}": "+Acme Corp"}}
        template_path: Path to template PPTX
        output_path: Path for output PPTX
        use_batch_replacer: If True, use python-pptx-text-replacer (recommended)

    Returns:
        Path to generated file
    """
    # Validate template exists
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(template_path)

    # Collect all replacements for batch processing
    all_replacements = {}

    for slide_index, replacements in content.items():
        if slide_index >= len(prs.slides):
            print(f"Warning: Slide index {slide_index} exceeds template size")
            continue

        for placeholder, text in replacements.items():
            all_replacements[placeholder] = text

    # Save intermediate file for batch replacer
    prs.save(output_path)

    # Use batch replacer if available and requested (preserves formatting best)
    if use_batch_replacer and HAS_TEXT_REPLACER:
        try:
            replace_text_batch(output_path, output_path, all_replacements)
            print(f"Used batch text replacer for {len(all_replacements)} replacements")
            return output_path
        except Exception as e:
            print(f"Batch replacer failed: {e}. Falling back to manual replacement.")

    # Fallback: manual replacement
    prs = Presentation(output_path)

    for slide_index, replacements in content.items():
        if slide_index >= len(prs.slides):
            continue

        slide = prs.slides[slide_index]
        for placeholder, text in replacements.items():
            if not replace_placeholder(slide, placeholder, text):
                print(f"Warning: Placeholder '{placeholder}' not found on slide {slide_index + 1}")

    prs.save(output_path)
    return output_path


def generate_presentation_with_duplicates(slide_specs, template_path=TEMPLATE_PATH,
                                           output_path=OUTPUT_PATH):
    """
    Generate a presentation by duplicating template slides and replacing content.

    This is the main workflow function for creating new presentations.

    Args:
        slide_specs: List of dicts, each with:
            - 'template_index': Which template slide to duplicate
            - 'replacements': Dict of {placeholder: text} pairs
        template_path: Path to template PPTX
        output_path: Path for output PPTX

    Returns:
        Path to generated file

    Example:
        slide_specs = [
            {
                'template_index': 0,  # Title slide
                'replacements': {
                    '{{title.client_name}}': '+ACME CORP',
                    '{{title.presentation_name}}': 'Q4 STRATEGY'
                }
            },
            {
                'template_index': 2,  # Section divider
                'replacements': {
                    '{{section.number}}': '01',
                    '{{section.title}}': 'THE CHALLENGE'
                }
            }
        ]
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(template_path)
    original_slide_count = len(prs.slides)

    # Create all duplicated slides first
    new_slides = []
    for spec in slide_specs:
        template_idx = spec.get('template_index', 0)
        if template_idx >= original_slide_count:
            print(f"Warning: Template index {template_idx} exceeds template size")
            continue

        new_slide = duplicate_slide(prs, template_idx)
        new_slides.append({
            'slide': new_slide,
            'replacements': spec.get('replacements', {})
        })

    # Remove original template slides (keep only duplicates)
    # We do this by saving and reloading with only the new slides
    # Actually, let's keep original slides and just work with new ones

    # Apply replacements to new slides
    for slide_info in new_slides:
        slide = slide_info['slide']
        for placeholder, text in slide_info['replacements'].items():
            if not replace_placeholder(slide, placeholder, text):
                print(f"Warning: Placeholder '{placeholder}' not found")

    prs.save(output_path)

    # Now use batch replacer for any remaining replacements
    if HAS_TEXT_REPLACER:
        all_replacements = {}
        for slide_info in new_slides:
            all_replacements.update(slide_info['replacements'])

        if all_replacements:
            try:
                replace_text_batch(output_path, output_path, all_replacements)
            except Exception as e:
                print(f"Note: Batch replacer skipped (already replaced manually): {e}")

    return output_path


# Example usage
if __name__ == "__main__":
    # Example content structure using the new workflow
    EXAMPLE_SPECS = [
        {
            'template_index': 0,  # Title slide
            'replacements': {
                '{{title.client_name}}': '+ACME CORP',
                '{{title.presentation_name}}': 'Q4 BRAND STRATEGY',
            }
        },
        {
            'template_index': 2,  # Section divider
            'replacements': {
                '{{section.number}}': '01',
                '{{section.title}}': 'THE CHALLENGE',
            }
        },
        {
            'template_index': 4,  # Statement dark
            'replacements': {
                '{{statement_dark.text}}': 'WE NEED TO REACH GEN Z WITHOUT ALIENATING OUR CORE.',
            }
        },
        {
            'template_index': 12,  # Statistics grid
            'replacements': {
                '{{stats.title}}': 'KEY METRICS',
                '{{stats.description}}': 'Our research revealed these critical insights:',
                '{{stats.value_1}}': '47%',
                '{{stats.label_1}}': 'Gen Z Discovery',
                '{{stats.desc_1}}': 'Find brands on TikTok',
                '{{stats.value_2}}': '3.2x',
                '{{stats.label_2}}': 'Engagement',
                '{{stats.desc_2}}': 'Higher with authentic content',
            }
        },
    ]

    print("HCO Presentation Generator V2")
    print("-" * 40)
    print(f"Template: {TEMPLATE_PATH}")
    print(f"Output: {OUTPUT_PATH}")
    print(f"Text Replacer Available: {HAS_TEXT_REPLACER}")
    print("-" * 40)

    try:
        # For simple replacement on existing slides:
        # result = generate_presentation(EXAMPLE_CONTENT)

        # For creating new presentations with duplicated slides:
        result = generate_presentation_with_duplicates(EXAMPLE_SPECS)

        print(f"\nGenerated: {result}")
        print("\nNote: This is example content. Modify EXAMPLE_SPECS")
        print("or import the functions for your own workflow.")
    except FileNotFoundError as e:
        print(f"\nError: {e}")
        print("Run from project root or update TEMPLATE_PATH")
        sys.exit(1)
    except Exception as e:
        print(f"\nError: {e}")
        sys.exit(1)
