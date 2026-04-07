#!/usr/bin/env python3
"""
Add 4 new fundamental slide templates to ppt-template-v2.pptx.

Creates:
- Slide 25 (index 24): bullet-list - Simple vertical bullet list
- Slide 26 (index 25): single-stat - Hero statistic with sidebar
- Slide 27 (index 26): heading-paragraph - Title + body text
- Slide 28 (index 27): highlights - Large callout statements

Run from skills/create-hco-ppt-v2 directory:
    python scripts/add_new_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from copy import deepcopy
import os

# Brand colors
CYAN = RGBColor(125, 211, 232)  # #7DD3E8
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

# Slide dimensions
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)


def set_line_spacing(paragraph, multiple):
    """Set line spacing as multiple (0.7 = 70%)."""
    pPr = paragraph._p.get_or_add_pPr()
    lnSpc = parse_xml(f'<a:lnSpc xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                      f'<a:spcPct val="{int(multiple * 100000)}"/></a:lnSpc>')
    # Remove existing lnSpc if present
    existing = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    if existing is not None:
        pPr.remove(existing)
    pPr.insert(0, lnSpc)


def add_text_box(slide, left, top, width, height, text, font_name, font_size,
                 font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a text box with specified formatting."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set anchor
    tf.paragraphs[0].alignment = alignment

    # Set anchor position (top, middle, bottom)
    textbox.text_frame._txBody.bodyPr.set('anchor',
        {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}[anchor])

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color

    if line_spacing:
        set_line_spacing(p, line_spacing)

    return textbox


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def copy_sidebar_from_slide(source_slide, target_slide):
    """Copy the sidebar group from a source slide to target slide."""
    for shape in source_slide.shapes:
        if hasattr(shape, 'shapes'):  # It's a group
            # Check if this looks like the sidebar (narrow width, full height)
            if shape.width < Inches(0.5) and shape.height > Inches(7):
                # This is likely the sidebar - we need to copy it
                # Clone the shape XML
                sp_tree = target_slide.shapes._spTree
                new_elem = deepcopy(shape._element)
                sp_tree.append(new_elem)
                return True
    return False


def create_bullet_list_slide(prs):
    """Create bullet-list slide (index 24) - Simple vertical bullet list."""
    # Use blank layout
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # White background (default)
    set_slide_background(slide, WHITE)

    # Title - VC Nudge Black, 48pt, UPPERCASE, 0.7 line spacing
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(0.58),
        width=Inches(12.15), height=Inches(1.0),
        text="{{bullet_list.title}}",
        font_name="VC Nudge Black", font_size=Pt(48),
        font_color=BLACK, line_spacing=0.7
    )

    # Items 1-7 - Manrope Light, 14pt, 0.60in vertical spacing
    items_start_top = 2.00
    item_spacing = 0.60

    for i in range(1, 8):
        add_text_box(
            slide,
            left=Inches(0.59), top=Inches(items_start_top + (i-1) * item_spacing),
            width=Inches(12.15), height=Inches(0.50),
            text=f"{{{{bullet_list.item_{i}}}}}",
            font_name="Manrope Light", font_size=Pt(14),
            font_color=BLACK
        )

    print("  Created slide 25: bullet-list")
    return slide


def create_single_stat_slide(prs, sidebar_source_slide):
    """Create single-stat slide (index 25) - Hero statistic with sidebar."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Black background
    set_slide_background(slide, BLACK)

    # Copy sidebar from statement-dark slide (index 4)
    if sidebar_source_slide:
        copy_sidebar_from_slide(sidebar_source_slide, slide)

    # Value - VC Nudge Black, 200pt, Cyan, centered, 0.7 line spacing
    add_text_box(
        slide,
        left=Inches(2.90), top=Inches(1.80),
        width=Inches(7.54), height=Inches(2.50),
        text="{{single_stat.value}}",
        font_name="VC Nudge Black", font_size=Pt(200),
        font_color=CYAN, alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.7
    )

    # Label - VC Nudge Black, 32pt, White, centered, 0.7 line spacing
    add_text_box(
        slide,
        left=Inches(2.90), top=Inches(4.50),
        width=Inches(7.54), height=Inches(0.80),
        text="{{single_stat.label}}",
        font_name="VC Nudge Black", font_size=Pt(32),
        font_color=WHITE, alignment=PP_ALIGN.CENTER,
        line_spacing=0.7
    )

    # Context - Manrope Light, 14pt, White, centered
    add_text_box(
        slide,
        left=Inches(2.90), top=Inches(5.50),
        width=Inches(7.54), height=Inches(0.50),
        text="{{single_stat.context}}",
        font_name="Manrope Light", font_size=Pt(14),
        font_color=WHITE, alignment=PP_ALIGN.CENTER
    )

    print("  Created slide 26: single-stat")
    return slide


def create_heading_paragraph_slide(prs):
    """Create heading-paragraph slide (index 26) - Title + body text."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # White background
    set_slide_background(slide, WHITE)

    # Heading - VC Nudge Black, 48pt, UPPERCASE, 0.7 line spacing, anchor bottom
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(1.50),
        width=Inches(12.15), height=Inches(1.83),
        text="{{heading_paragraph.heading}}",
        font_name="VC Nudge Black", font_size=Pt(48),
        font_color=BLACK, anchor=MSO_ANCHOR.BOTTOM,
        line_spacing=0.7
    )

    # Body - Manrope Light, 14pt
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(3.75),
        width=Inches(12.15), height=Inches(3.00),
        text="{{heading_paragraph.body}}",
        font_name="Manrope Light", font_size=Pt(14),
        font_color=BLACK
    )

    print("  Created slide 27: heading-paragraph")
    return slide


def create_highlights_slide(prs):
    """Create highlights slide (index 27) - Large callout statements."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # White background
    set_slide_background(slide, WHITE)

    # Title - VC Nudge Black, 48pt, UPPERCASE, 0.7 line spacing
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(0.58),
        width=Inches(12.15), height=Inches(1.0),
        text="{{highlights.title}}",
        font_name="VC Nudge Black", font_size=Pt(48),
        font_color=BLACK, line_spacing=0.7
    )

    # Items 1-5 - Manrope Light, 24pt, 0.90in vertical spacing
    items_start_top = 2.00
    item_spacing = 0.90

    for i in range(1, 6):
        add_text_box(
            slide,
            left=Inches(0.59), top=Inches(items_start_top + (i-1) * item_spacing),
            width=Inches(12.15), height=Inches(0.70),
            text=f"{{{{highlights.item_{i}}}}}",
            font_name="Manrope Light", font_size=Pt(24),
            font_color=BLACK
        )

    print("  Created slide 28: highlights")
    return slide


def main():
    # Paths
    script_dir = os.path.dirname(os.path.abspath(__file__))
    skill_dir = os.path.dirname(script_dir)
    template_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")
    output_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")

    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)

    print(f"Current slide count: {len(prs.slides)}")

    # Get sidebar source (slide 5, index 4 - statement-dark)
    sidebar_source = prs.slides[4] if len(prs.slides) > 4 else None

    print("\nCreating new slides...")

    # Create the 4 new slides
    create_bullet_list_slide(prs)
    create_single_stat_slide(prs, sidebar_source)
    create_heading_paragraph_slide(prs)
    create_highlights_slide(prs)

    print(f"\nNew slide count: {len(prs.slides)}")

    # Save
    print(f"\nSaving to: {output_path}")
    prs.save(output_path)
    print("Done!")


if __name__ == "__main__":
    main()
