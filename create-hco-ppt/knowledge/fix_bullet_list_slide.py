#!/usr/bin/env python3
"""
Fix bullet-list slide (Index 24) to use single text box with proper bullet formatting.

The original slide used 7 separate text boxes for each item. This script replaces
it with a single text box containing 7 paragraphs with bullet characters.

Run from skills/create-hco-ppt-v2 directory:
    python scripts/fix_bullet_list_slide.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Brand colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)


def set_line_spacing(paragraph, multiple):
    """Set line spacing as multiple (0.7 = 70%)."""
    pPr = paragraph._p.get_or_add_pPr()
    lnSpc = parse_xml(f'<a:lnSpc xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                      f'<a:spcPct val="{int(multiple * 100000)}"/></a:lnSpc>')
    existing = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    if existing is not None:
        pPr.remove(existing)
    pPr.insert(0, lnSpc)


def add_bullet_char(paragraph):
    """Add bullet character formatting to a paragraph."""
    pPr = paragraph._p.get_or_add_pPr()
    # Add bullet character
    buChar = parse_xml('<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="•"/>')
    # Remove existing bullet if present
    existing = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
    if existing is not None:
        pPr.remove(existing)
    # Insert after buNone if it exists, otherwise at beginning
    buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
    if buNone is not None:
        pPr.remove(buNone)
    pPr.append(buChar)


def set_paragraph_spacing(paragraph, before_pt=0, after_pt=12):
    """Set paragraph spacing before and after in points."""
    pPr = paragraph._p.get_or_add_pPr()
    # Space before
    spcBef = parse_xml(f'<a:spcBef xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                       f'<a:spcPts val="{int(before_pt * 100)}"/></a:spcBef>')
    existing = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spcBef')
    if existing is not None:
        pPr.remove(existing)
    pPr.append(spcBef)
    # Space after
    spcAft = parse_xml(f'<a:spcAft xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                       f'<a:spcPts val="{int(after_pt * 100)}"/></a:spcAft>')
    existing = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft')
    if existing is not None:
        pPr.remove(existing)
    pPr.append(spcAft)


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_bullet_list_slide(prs, insert_index):
    """Create bullet-list slide with single text box containing bulleted paragraphs."""
    # Use blank layout
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # White background
    set_slide_background(slide, WHITE)

    # Title - VC Nudge Black, 48pt, UPPERCASE, 0.7 line spacing
    title_box = slide.shapes.add_textbox(
        Inches(0.59), Inches(0.58),
        Inches(12.15), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    title_box.text_frame._txBody.bodyPr.set('anchor', 't')

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{bullet_list.title}}"
    run.font.name = "VC Nudge Black"
    run.font.size = Pt(48)
    run.font.bold = False
    run.font.color.rgb = BLACK
    set_line_spacing(p, 0.7)

    # Bullet list - single text box with 7 bulleted paragraphs
    bullet_box = slide.shapes.add_textbox(
        Inches(0.59), Inches(2.00),
        Inches(12.15), Inches(4.75)
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bullet_box.text_frame._txBody.bodyPr.set('anchor', 't')

    # First paragraph
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "{{bullet_list.item_1}}"
    run.font.name = "Manrope Light"
    run.font.size = Pt(14)
    run.font.color.rgb = BLACK
    add_bullet_char(p)
    set_paragraph_spacing(p, before_pt=0, after_pt=12)

    # Add remaining 6 paragraphs
    for i in range(2, 8):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{{{{bullet_list.item_{i}}}}}"
        run.font.name = "Manrope Light"
        run.font.size = Pt(14)
        run.font.color.rgb = BLACK
        add_bullet_char(p)
        set_paragraph_spacing(p, before_pt=0, after_pt=12)

    print("  Created new bullet-list slide with proper bullet formatting")
    return slide


def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index."""
    slides = list(prs.slides._sldIdLst)
    slide_to_move = slides.pop(old_index)
    slides.insert(new_index, slide_to_move)
    prs.slides._sldIdLst[:] = slides


def delete_slide(prs, index):
    """Delete a slide at the given index."""
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    skill_dir = os.path.dirname(script_dir)
    template_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")
    output_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")

    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)

    initial_count = len(prs.slides)
    print(f"Current slide count: {initial_count}")

    # Index 24 is the bullet-list slide
    target_index = 24

    print(f"\nStep 1: Deleting old bullet-list slide at index {target_index}")
    delete_slide(prs, target_index)
    print(f"  Slide count after deletion: {len(prs.slides)}")

    print(f"\nStep 2: Creating new bullet-list slide")
    create_bullet_list_slide(prs, target_index)
    print(f"  Slide count after creation: {len(prs.slides)}")

    # The new slide is added at the end, move it to index 24
    print(f"\nStep 3: Moving new slide from end to index {target_index}")
    move_slide(prs, len(prs.slides) - 1, target_index)

    print(f"\nFinal slide count: {len(prs.slides)}")

    print(f"\nSaving to: {output_path}")
    prs.save(output_path)
    print("Done!")


if __name__ == "__main__":
    main()
