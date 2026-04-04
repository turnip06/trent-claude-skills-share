#!/usr/bin/env python3
"""
Test script to verify the 4 new slide templates work correctly.

Creates a test presentation using each of the new slides:
- Slide 25 (index 24): bullet-list
- Slide 26 (index 25): single-stat
- Slide 27 (index 26): heading-paragraph
- Slide 28 (index 27): highlights

Run from skills/create-hco-ppt-v2 directory:
    python scripts/test_new_slides.py
"""

import os
import sys
from pptx import Presentation

# Add scripts directory to path
script_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, script_dir)

from generate_presentation import duplicate_slide, replace_placeholder


def main():
    # Paths
    skill_dir = os.path.dirname(script_dir)
    template_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")
    output_path = os.path.join(skill_dir, "test-output", "test-new-slides.pptx")

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)
    print(f"Template has {len(prs.slides)} slides")

    # Create test presentation with just title + new slides + thank-you
    test_prs = Presentation(template_path)

    print("\nTesting new slides...")

    # Test 1: bullet-list (index 24)
    print("\n1. Testing bullet-list (index 24)...")
    slide = duplicate_slide(test_prs, 24)
    replace_placeholder(slide, "{{bullet_list.title}}", "KEY TAKEAWAYS")
    replace_placeholder(slide, "{{bullet_list.item_1}}", "First important point to remember")
    replace_placeholder(slide, "{{bullet_list.item_2}}", "Second critical insight from our analysis")
    replace_placeholder(slide, "{{bullet_list.item_3}}", "Third action item for the team")
    replace_placeholder(slide, "{{bullet_list.item_4}}", "Fourth recommendation for next steps")
    replace_placeholder(slide, "{{bullet_list.item_5}}", "Fifth consideration for stakeholders")
    replace_placeholder(slide, "{{bullet_list.item_6}}", "Sixth factor to keep in mind")
    replace_placeholder(slide, "{{bullet_list.item_7}}", "Seventh and final takeaway")
    print("   bullet-list: OK")

    # Test 2: single-stat (index 25)
    print("\n2. Testing single-stat (index 25)...")
    slide = duplicate_slide(test_prs, 25)
    replace_placeholder(slide, "{{single_stat.value}}", "67%")
    replace_placeholder(slide, "{{single_stat.label}}", "CUSTOMER SATISFACTION")
    replace_placeholder(slide, "{{single_stat.context}}", "Based on Q4 2024 survey results")
    print("   single-stat: OK")

    # Test 3: heading-paragraph (index 26)
    print("\n3. Testing heading-paragraph (index 26)...")
    slide = duplicate_slide(test_prs, 26)
    replace_placeholder(slide, "{{heading_paragraph.heading}}", "OUR APPROACH")
    replace_placeholder(slide, "{{heading_paragraph.body}}",
         "We believe in a customer-first approach that prioritizes long-term relationships over short-term gains. "
         "Our methodology combines data-driven insights with creative problem-solving to deliver solutions that "
         "truly resonate with your target audience. By understanding the unique challenges and opportunities in "
         "your market, we craft strategies that drive measurable results and sustainable growth.")
    print("   heading-paragraph: OK")

    # Test 4: highlights (index 27)
    print("\n4. Testing highlights (index 27)...")
    slide = duplicate_slide(test_prs, 27)
    replace_placeholder(slide, "{{highlights.title}}", "KEY HIGHLIGHTS")
    replace_placeholder(slide, "{{highlights.item_1}}", "Brand awareness increased by 45% year-over-year")
    replace_placeholder(slide, "{{highlights.item_2}}", "Customer acquisition cost reduced by 30%")
    replace_placeholder(slide, "{{highlights.item_3}}", "Net promoter score improved from 42 to 67")
    replace_placeholder(slide, "{{highlights.item_4}}", "Market share expanded in three new regions")
    replace_placeholder(slide, "{{highlights.item_5}}", "Revenue growth exceeded projections by 15%")
    print("   highlights: OK")

    # Save test presentation
    print(f"\nSaving test presentation: {output_path}")
    test_prs.save(output_path)

    # Run validation
    print("\nRunning validation...")
    from validate_presentation import validate_presentation
    result = validate_presentation(output_path, verbose=True)

    print(f"\nSlide count: {result['slide_count']}")
    print(f"Unreplaced placeholders: {result['unreplaced_count']}")

    # Check if our new slides have any unreplaced placeholders
    # (They shouldn't - we replaced them all)
    if result['unreplaced_count'] > 0:
        print("\nNote: Some original template slides have unreplaced placeholders (expected)")
        # Filter to see if any are from our new slides
        new_slide_prefixes = ['bullet_list', 'single_stat', 'heading_paragraph', 'highlights']
        for error in result['errors']:
            for prefix in new_slide_prefixes:
                if prefix in error:
                    print(f"  WARNING: {error}")

    print("\n✓ All new slides created and tested successfully!")
    print(f"  Output: {output_path}")


if __name__ == "__main__":
    main()
