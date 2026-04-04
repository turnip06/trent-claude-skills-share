#!/usr/bin/env python3
"""
Add 3 new slide templates to ppt-template-v2.pptx.

Creates:
- Slide 29 (index 28): horizontal-split - Two horizontal sections (cyan/white)
- Slide 30 (index 29): split-title-content - Vertical split with title on left
- Slide 31 (index 30): final-thought - Black bg with cyan heading

Run from skills/create-hco-ppt-v2 directory:
    python scripts/add_template_slides_v3.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from copy import deepcopy
import os

# Brand colors
CYAN = RGBColor(125, 211, 232)  # #7DD3E8
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

# Slide dimensions
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)


def set_line_spacing(paragraph, multiple):
    """Set line spacing as multiple (0.7 = 70%)."""
    pPr = paragraph._p.get_or_add_pPr()
    lnSpc = parse_xml(f'<a:lnSpc xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                      f'<a:spcPct val="{int(multiple * 100000)}"/></a:lnSpc>')
    existing = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc')
    if existing is not None:
        pPr.remove(existing)
    pPr.insert(0, lnSpc)


def add_text_box(slide, left, top, width, height, text, font_name, font_size,
                 font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Add a text box with specified formatting."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    tf.paragraphs[0].alignment = alignment
    textbox.text_frame._txBody.bodyPr.set('anchor',
        {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}[anchor])

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = font_color

    if line_spacing:
        set_line_spacing(p, line_spacing)

    return textbox


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def copy_sidebar_from_slide(source_slide, target_slide):
    """Copy the sidebar group from a source slide to target slide."""
    for shape in source_slide.shapes:
        if hasattr(shape, 'shapes'):  # It's a group
            if shape.width < Inches(0.5) and shape.height > Inches(7):
                sp_tree = target_slide.shapes._spTree
                new_elem = deepcopy(shape._element)
                sp_tree.append(new_elem)
                return True
    return False


def create_horizontal_split_slide(prs):
    """Create horizontal-split slide (index 28) - Two horizontal sections."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # White background (bottom half default)
    set_slide_background(slide, WHITE)

    # Cyan rectangle for top half
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.75), CYAN)

    # TOP SECTION (Cyan background)
    # Title
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(0.40),
        width=Inches(5.5), height=Inches(0.60),
        text="{{horizontal.title_top}}",
        font_name="VC Nudge Black", font_size=Pt(28),
        font_color=BLACK, line_spacing=0.7
    )

    # Subtitle
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(1.00),
        width=Inches(5.5), height=Inches(0.40),
        text="{{horizontal.subtitle_top}}",
        font_name="Manrope Light", font_size=Pt(12),
        font_color=BLACK
    )

    # Column 1a - heading
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(1.60),
        width=Inches(5.5), height=Inches(0.40),
        text="{{horizontal.head_1a}}",
        font_name="Manrope", font_size=Pt(14),
        font_color=BLACK, bold=True
    )

    # Column 1a - body
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(2.00),
        width=Inches(5.5), height=Inches(1.50),
        text="{{horizontal.body_1a}}",
        font_name="Manrope Light", font_size=Pt(11),
        font_color=BLACK
    )

    # Column 1b - heading
    add_text_box(
        slide,
        left=Inches(6.70), top=Inches(1.60),
        width=Inches(5.5), height=Inches(0.40),
        text="{{horizontal.head_1b}}",
        font_name="Manrope", font_size=Pt(14),
        font_color=BLACK, bold=True
    )

    # Column 1b - body
    add_text_box(
        slide,
        left=Inches(6.70), top=Inches(2.00),
        width=Inches(5.5), height=Inches(1.50),
        text="{{horizontal.body_1b}}",
        font_name="Manrope Light", font_size=Pt(11),
        font_color=BLACK
    )

    # BOTTOM SECTION (White background)
    # Title
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(4.00),
        width=Inches(5.5), height=Inches(0.60),
        text="{{horizontal.title_bottom}}",
        font_name="VC Nudge Black", font_size=Pt(28),
        font_color=BLACK, line_spacing=0.7
    )

    # Subtitle
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(4.60),
        width=Inches(5.5), height=Inches(0.40),
        text="{{horizontal.subtitle_bottom}}",
        font_name="Manrope Light", font_size=Pt(12),
        font_color=BLACK
    )

    # Column 2a - heading
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(5.20),
        width=Inches(5.5), height=Inches(0.40),
        text="{{horizontal.head_2a}}",
        font_name="Manrope", font_size=Pt(14),
        font_color=BLACK, bold=True
    )

    # Column 2a - body
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(5.60),
        width=Inches(5.5), height=Inches(1.50),
        text="{{horizontal.body_2a}}",
        font_name="Manrope Light", font_size=Pt(11),
        font_color=BLACK
    )

    # Column 2b - heading
    add_text_box(
        slide,
        left=Inches(6.70), top=Inches(5.20),
        width=Inches(5.5), height=Inches(0.40),
        text="{{horizontal.head_2b}}",
        font_name="Manrope", font_size=Pt(14),
        font_color=BLACK, bold=True
    )

    # Column 2b - body
    add_text_box(
        slide,
        left=Inches(6.70), top=Inches(5.60),
        width=Inches(5.5), height=Inches(1.50),
        text="{{horizontal.body_2b}}",
        font_name="Manrope Light", font_size=Pt(11),
        font_color=BLACK
    )

    print("  Created slide 29: horizontal-split")
    return slide


def create_split_title_content_slide(prs):
    """Create split-title-content slide (index 29) - Title on left, content on right."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # White background (default)
    set_slide_background(slide, WHITE)

    # Cyan rectangle for left half
    add_rectangle(slide, Inches(0), Inches(0), Inches(6.67), SLIDE_HEIGHT, CYAN)

    # LEFT SIDE (Cyan background)
    # Title - centered vertically in left half
    add_text_box(
        slide,
        left=Inches(0.59), top=Inches(2.50),
        width=Inches(5.50), height=Inches(2.50),
        text="{{split.title}}",
        font_name="VC Nudge Black", font_size=Pt(48),
        font_color=BLACK, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=0.7
    )

    # RIGHT SIDE (White background)
    # Heading 1
    add_text_box(
        slide,
        left=Inches(7.25), top=Inches(1.50),
        width=Inches(5.50), height=Inches(0.50),
        text="{{split.head_1}}",
        font_name="Manrope", font_size=Pt(14),
        font_color=BLACK, bold=True
    )

    # Body 1
    add_text_box(
        slide,
        left=Inches(7.25), top=Inches(2.00),
        width=Inches(5.50), height=Inches(1.80),
        text="{{split.body_1}}",
        font_name="Manrope Light", font_size=Pt(11),
        font_color=BLACK
    )

    # Heading 2
    add_text_box(
        slide,
        left=Inches(7.25), top=Inches(4.20),
        width=Inches(5.50), height=Inches(0.50),
        text="{{split.head_2}}",
        font_name="Manrope", font_size=Pt(14),
        font_color=BLACK, bold=True
    )

    # Body 2
    add_text_box(
        slide,
        left=Inches(7.25), top=Inches(4.70),
        width=Inches(5.50), height=Inches(1.80),
        text="{{split.body_2}}",
        font_name="Manrope Light", font_size=Pt(11),
        font_color=BLACK
    )

    print("  Created slide 30: split-title-content")
    return slide


def create_final_thought_slide(prs, sidebar_source_slide):
    """Create final-thought slide (index 30) - Black bg with cyan heading."""
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Black background
    set_slide_background(slide, BLACK)

    # Copy sidebar from statement-dark slide
    if sidebar_source_slide:
        copy_sidebar_from_slide(sidebar_source_slide, slide)

    # Heading - Cyan, VC Nudge Black, 48pt
    add_text_box(
        slide,
        left=Inches(1.50), top=Inches(2.20),
        width=Inches(10.50), height=Inches(1.50),
        text="{{final.heading}}",
        font_name="VC Nudge Black", font_size=Pt(48),
        font_color=CYAN, line_spacing=0.7
    )

    # Body - White, Manrope Light, 16pt
    add_text_box(
        slide,
        left=Inches(1.50), top=Inches(4.00),
        width=Inches(10.50), height=Inches(2.00),
        text="{{final.body}}",
        font_name="Manrope Light", font_size=Pt(16),
        font_color=WHITE
    )

    print("  Created slide 31: final-thought")
    return slide


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    skill_dir = os.path.dirname(script_dir)
    template_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")
    output_path = os.path.join(skill_dir, "assets", "ppt-template-v2.pptx")

    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)

    print(f"Current slide count: {len(prs.slides)}")

    # Get sidebar source (slide 5, index 4 - statement-dark)
    sidebar_source = prs.slides[4] if len(prs.slides) > 4 else None

    print("\nCreating new slides...")

    create_horizontal_split_slide(prs)
    create_split_title_content_slide(prs)
    create_final_thought_slide(prs, sidebar_source)

    print(f"\nNew slide count: {len(prs.slides)}")

    print(f"\nSaving to: {output_path}")
    prs.save(output_path)
    print("Done!")


if __name__ == "__main__":
    main()
