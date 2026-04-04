#!/usr/bin/env python3
"""
Generate Grime-X Creative Pitch presentation.
25 slides mapped to HCO template.

V3 UPDATES:
- Uses improved duplicate_slide() that preserves backgrounds and fixes image rIds
- Applies heading height adjustments for long text
- Applies section title position adjustments for wide numbers
"""

import sys
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import improved functions from generate_presentation
from scripts.generate_presentation import (
    duplicate_slide,
    adjust_heading_height,
    adjust_section_title_position
)

# Try to import the text replacer package
try:
    from python_pptx_text_replacer import TextReplacer
    HAS_TEXT_REPLACER = True
except ImportError:
    HAS_TEXT_REPLACER = False

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                              "assets/ppt-template-v2.pptx")
OUTPUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                            "output.pptx")


def replace_placeholder(slide, placeholder, new_text):
    """Replace placeholder text in all shapes on a slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def replace_in_shape(shape):
        replaced = False
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full_text = ''.join(run.text for run in para.runs)
                if placeholder in full_text:
                    for run in para.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, new_text)
                            replaced = True
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if replace_in_shape(child):
                    replaced = True
        return replaced

    replaced = False
    for shape in slide.shapes:
        if replace_in_shape(shape):
            replaced = True
    return replaced


def generate_presentation_clean(slide_specs, template_path, output_path):
    """Generate presentation keeping only duplicated slides (removing template slides)."""
    prs = Presentation(template_path)
    original_count = len(prs.slides)

    # Slide types that need heading height adjustment (text-image-right = index 15)
    HEADING_ADJUSTMENT_SLIDES = {15}
    # Slide types that need section title adjustment (section-divider = index 2)
    SECTION_DIVIDER_INDEX = 2

    # First, duplicate all slides with their content
    new_slide_indices = []
    for spec in slide_specs:
        template_idx = spec.get('template_index', 0)
        new_slide = duplicate_slide(prs, template_idx)
        new_slide_indices.append(prs.slides.index(new_slide))

        # Apply adjustments based on slide type
        replacements = spec.get('replacements', {})

        # Adjust heading height for text-image slides with long content
        if template_idx in HEADING_ADJUSTMENT_SLIDES:
            heading_text = replacements.get('{{text_image.heading}}', '')
            if heading_text:
                for shape in new_slide.shapes:
                    if shape.has_text_frame and '{{text_image.heading}}' in shape.text_frame.text:
                        adjust_heading_height(shape, heading_text, max_chars_per_line=25)
                        break

        # Adjust section title position for wide numbers
        if template_idx == SECTION_DIVIDER_INDEX:
            section_number = replacements.get('{{section.number}}', '01')
            adjust_section_title_position(new_slide, section_number)

        # Apply replacements
        for placeholder, text in replacements.items():
            if not replace_placeholder(new_slide, placeholder, text):
                print(f"Warning: Placeholder '{placeholder}' not found")

    # Remove original template slides (indices 0 to original_count-1)
    # We need to remove from highest index to lowest to preserve indices
    for i in range(original_count - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    prs.save(output_path)

    # Run batch text replacer if available for any remaining placeholders
    if HAS_TEXT_REPLACER:
        all_replacements = {}
        for spec in slide_specs:
            all_replacements.update(spec.get('replacements', {}))
        if all_replacements:
            try:
                replacer = TextReplacer(output_path, tables=True, charts=True, textframes=True)
                replacement_list = [(k, v) for k, v in all_replacements.items()]
                replacer.replace_text(replacement_list)
                replacer.write_presentation_to_file(output_path)
            except:
                pass

    return output_path

# All 25 slides
SLIDE_SPECS = [
    # Slide 1: Title
    {
        'template_index': 0,
        'replacements': {
            '{{title.client_name}}': '+GRIME-X',
            '{{title.presentation_name}}': 'CREATIVE PITCH',
        }
    },
    # Slide 2: Statement Dark - Hygiene Boom
    {
        'template_index': 4,
        'replacements': {
            '{{statement_dark.text}}': "WE'RE HERE TO TALK ABOUT THE HYGIENE BOOM",
        }
    },
    # Slide 3: Text + Image - Product intro
    {
        'template_index': 15,
        'replacements': {
            '{{text_image.heading}}': "MORE SPECIFICALLY, WHAT WE'RE GOING TO DO ABOUT IT... NATURALLY.",
            '{{text_image.body}}': 'A bottle of Grime-X Plant-Based Heavy Duty Spray.',
        }
    },
    # Slide 4: Solid Bubbles - Category pressure
    {
        'template_index': 11,
        'replacements': {
            '{{solid_bubbles.preheader}}': 'WE KNOW THIS CATEGORY',
            '{{solid_bubbles.title}}': 'LIVES WITH PRESSURE FROM ALL SIDES',
            '{{solid_bubbles.bubble_1}}': 'CONSUMERS',
            '{{solid_bubbles.body_1}}': 'Demanding safer products',
            '{{solid_bubbles.bubble_2}}': 'ENVIRONMENTAL GROUPS',
            '{{solid_bubbles.body_2}}': 'Pushing for sustainability',
            '{{solid_bubbles.bubble_3}}': 'HEALTH REGULATORS',
            '{{solid_bubbles.body_3}}': 'Tightening standards',
        }
    },
    # Slide 5: Statement Light - Make people pay attention
    {
        'template_index': 5,
        'replacements': {
            '{{statement_light.text}}': 'BUT WE ALSO KNOW THAT WE HAVE TO MAKE PEOPLE PAY ATTENTION.',
        }
    },
    # Slide 6: Solid Bubbles - Today agenda
    {
        'template_index': 11,
        'replacements': {
            '{{solid_bubbles.preheader}}': 'TODAY',
            '{{solid_bubbles.title}}': 'AGENDA',
            '{{solid_bubbles.bubble_1}}': 'STRATEGY RECAP',
            '{{solid_bubbles.body_1}}': 'Where we are',
            '{{solid_bubbles.bubble_2}}': 'CREATIVE CONCEPTS',
            '{{solid_bubbles.body_2}}': 'Five ideas to share',
            '{{solid_bubbles.bubble_3}}': 'SUMMARY',
            '{{solid_bubbles.body_3}}': 'Next steps',
        }
    },
    # Slide 7: Statement Dark - Objective
    {
        'template_index': 4,
        'replacements': {
            '{{statement_dark.text}}': 'OBJECTIVE: SELECT ONE CREATIVE CONCEPT TO DEVELOP FURTHER WITH ALL PARTNER AGENCIES',
        }
    },
    # Slide 8: Section Divider - Strategy
    {
        'template_index': 2,
        'replacements': {
            '{{section.number}}': '01',
            '{{section.title}}': 'STRATEGY',
        }
    },
    # Slide 9: Text + Image - Market leader
    {
        'template_index': 15,
        'replacements': {
            '{{text_image.heading}}': "WE'RE THE MARKET LEADER, BUT THE CATEGORY IS UNDER INCREASING PRESSURE",
            '{{text_image.body}}': 'Montage of news articles: "Study links harsh bleach to lung issues", "Parents ditching chemical cleaners for home remedies", "Chlorine banned in several school districts".',
        }
    },
    # Slide 10: Text + Image - Consumer sentiment
    {
        'template_index': 15,
        'replacements': {
            '{{text_image.heading}}': 'CONSUMERS ARE TURNING TO SOLUTIONS THAT ARE SAFE FOR THEIR FAMILIES, THEIR PETS AND THE ENVIRONMENT',
            '{{text_image.body}}': 'Montage of forum posts: "Common cleaners are irritating my dog\'s paws", "Safe heavy duty cleaner?" (Online Forum), "Mould in the shower again. Fume-free solutions?".',
        }
    },
    # Slide 11: Three Column Bubbles - Challenge data
    {
        'template_index': 10,
        'replacements': {
            '{{three_col_bubbles.title}}': 'THIS IS REFLECTED BY THE CHALLENGE WE IDENTIFIED IN THE BRIEF',
            '{{three_col_bubbles.bubble_1a}}': '-8%',
            '{{three_col_bubbles.bubble_1b}}': 'SAFE',
            '{{three_col_bubbles.body_1}}': 'Safe around children',
            '{{three_col_bubbles.bubble_2a}}': '-10%',
            '{{three_col_bubbles.bubble_2b}}': 'PETS',
            '{{three_col_bubbles.body_2}}': 'Safe around pets',
            '{{three_col_bubbles.bubble_3a}}': '-13%',
            '{{three_col_bubbles.bubble_3b}}': 'NATURAL',
            '{{three_col_bubbles.body_3}}': 'Effective Natural Solutions',
        }
    },
    # Slide 12: Statement Dark - Clear way in
    {
        'template_index': 4,
        'replacements': {
            '{{statement_dark.text}}': 'WHICH GAVE US A CLEAR WAY IN, ALL ANCHORED IN OUR NAME: OWN THE "X" BY PROVING IT\'S JUST AS EFFECTIVE, WHILE BEING SAFER.',
        }
    },
    # Slide 13: Statement Light - Nature is awesome
    {
        'template_index': 5,
        'replacements': {
            '{{statement_light.text}}': "BECAUSE NATURE IS POWERFUL, REAL, REJUVENATING AND PROTECTIVE... RUTHLESS. WHEN IT NEEDS TO BE.",
        }
    },
    # Slide 14: Three Column Text - Strategy + RTBs
    {
        'template_index': 6,
        'replacements': {
            '{{three_col_text.title}}': "THE STRATEGY: 'NO COMPROMISE'",
            '{{three_col_text.heading_1}}': 'RTB 1',
            '{{three_col_text.body_1}}': 'Based upon Natural active enzymes that contain no bleach and no ammonia or harsh fumes.',
            '{{three_col_text.heading_2}}': 'RTB 2',
            '{{three_col_text.body_2}}': 'Safe for skin and lungs. Does not irritate eyes, throat, or pets.',
            '{{three_col_text.heading_3}}': 'RTB 3',
            '{{three_col_text.body_3}}': 'Eliminates grease and mould and is available in spray and concentrate.',
        }
    },
    # Slide 15: Statement Dark - Ruthless insight
    {
        'template_index': 4,
        'replacements': {
            '{{statement_dark.text}}': "DURING OUR CREATIVE CONCEPTING, WE REALLY ATTACHED OURSELVES TO ONE PIECE: 'RUTHLESS'",
        }
    },
    # Slide 16: Fullbleed Image - Website screenshot
    {
        'template_index': 17,
        'replacements': {
            '{{fullbleed.text}}': "DON'T LET DIRTY NO GOOD FILTH INVADE YOUR HOME!",
        }
    },
    # Slide 17: Two Column Text - Creative ask
    {
        'template_index': 8,
        'replacements': {
            '{{two_col_text.title}}': 'THIS LED US TO A CREATIVE ASK IN TANDEM WITH OUR STRATEGIC PROPOSITION',
            '{{two_col_text.heading_1}}': 'STRATEGIC PROPOSITION',
            '{{two_col_text.body_1}}': 'NATURAL IS NEVER A COMPROMISE',
            '{{two_col_text.heading_2}}': 'CREATIVE LINE',
            '{{two_col_text.body_2}}': 'BAD FOR GRIME. NOT BAD FOR YOU.',
        }
    },
    # Slide 18: Section Divider - Creative
    {
        'template_index': 2,
        'replacements': {
            '{{section.number}}': '02',
            '{{section.title}}': 'CREATIVE',
        }
    },
    # Slide 19: Statement Dark - Five ideas
    {
        'template_index': 4,
        'replacements': {
            '{{statement_dark.text}}': 'WE HAVE FIVE IDEAS TO SHARE TODAY',
        }
    },
    # Slide 20: Timeline Process - 5 pillars
    {
        'template_index': 22,
        'replacements': {
            '{{timeline.header}}': 'EACH WITH A UNIQUE WAY IN TO OUR STRATEGY',
            '{{timeline.image}}': '',  # Will keep template image
            # Step 1 row (4 items)
            '{{timeline.num_1}}': '01',
            '{{timeline.step_1a}}': 'FULL VOLUME',
            '{{timeline.step_1b}}': 'Tackles the problem head on',
            '{{timeline.step_1c}}': '',
            '{{timeline.step_1d}}': '',
            # Step 2
            '{{timeline.num_2}}': '02',
            '{{timeline.step_2}}': 'BYPRODUCT',
            '{{timeline.step_2a}}': 'Focus on the effectiveness byproduct',
            # Step 3
            '{{timeline.num_3}}': '03',
            '{{timeline.step_3}}': 'BIG SOLUTION',
            '{{timeline.step_3a}}': 'As big as the problem we face',
            # Step 4
            '{{timeline.num_4}}': '04',
            '{{timeline.step_4a}}': 'BEHAVIOURAL',
            '{{timeline.step_4b}}': 'A shift to a reflex we all have',
            # Step 5
            '{{timeline.num_5}}': '05',
            '{{timeline.step_5a}}': 'PARTNERSHIP',
            '{{timeline.step_5b}}': 'Part celebration',
            '{{timeline.step_5c}}': 'Part intervention',
            # Step 6 (empty - only 5 pillars)
            '{{timeline.num_6}}': '',
            '{{timeline.step_6a}}': '',
            '{{timeline.step_6b}}': '',
            '{{timeline.step_7}}': '',
            # Callouts
            '{{timeline.callout_1}}': '> NATURAL IS NEVER A COMPROMISE',
            '{{timeline.callout_2}}': '> BAD FOR GRIME. NOT BAD FOR YOU.',
            '{{timeline.callout_3}}': '',
            '{{timeline.callout_4}}': '',
            '{{timeline.callout_5}}': '',
            '{{timeline.callout_6}}': '',
            '{{timeline.callout_7}}': '',
            '{{timeline.callout_8}}': '',
        }
    },
    # Slide 21: Text + Image - Provocation focus
    {
        'template_index': 15,
        'replacements': {
            '{{text_image.heading}}': "FOR EACH, WE'VE FOCUSED ON PROVOCATION. THIS IS WHAT WILL DRIVE OUTSIZED IMPACT.",
            '{{text_image.body}}': "BUT WE'RE NOT BEING PROVOCATIVE TOWARDS THEM. THIS IS ABOUT DIRT.\n\nLogos: Earth Alliance, National HBoard",
        }
    },
    # Slide 22: Statement Dark - Provocative
    {
        'template_index': 4,
        'replacements': {
            '{{statement_dark.text}}': "IF WE'RE NOT PROVOCATIVE NO ONE WILL GIVE A DAMN",
        }
    },
    # Slide 23: Statement Light - Love clean hate fumes
    {
        'template_index': 5,
        'replacements': {
            '{{statement_light.text}}': "SO IT'S FOR PEOPLE WHO LOVE CLEAN BUT HATE FUMES.",
        }
    },
    # Slide 24: Section Divider - Concept 01
    {
        'template_index': 2,
        'replacements': {
            '{{section.number}}': '01',
            '{{section.title}}': 'CONCEPT ONE',
        }
    },
    # Slide 25: Thank You
    {
        'template_index': 23,
        'replacements': {}  # Static slide, no replacements
    },
]

if __name__ == "__main__":
    print("Generating Grime-X Creative Pitch")
    print("-" * 40)
    print(f"Template: {TEMPLATE_PATH}")
    print(f"Output: {OUTPUT_PATH}")
    print(f"Slides: {len(SLIDE_SPECS)}")
    print(f"Text Replacer: {HAS_TEXT_REPLACER}")
    print("-" * 40)

    try:
        result = generate_presentation_clean(
            SLIDE_SPECS,
            template_path=TEMPLATE_PATH,
            output_path=OUTPUT_PATH
        )
        print(f"\nCreated: {result}")
        print(f"Generated {len(SLIDE_SPECS)} slides successfully!")
    except Exception as e:
        print(f"\nError: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
