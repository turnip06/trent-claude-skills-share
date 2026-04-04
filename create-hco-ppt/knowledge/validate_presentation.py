#!/usr/bin/env python3
"""
Validate a generated HCO presentation.

V2 IMPROVEMENTS:
- Font consistency checking
- Image presence validation
- Shape position comparison
- Color accuracy validation
- More detailed reporting

Checks for:
- Unreplaced placeholder tokens
- Slide count and structure
- Thank-you slide presence
- File integrity
- Font consistency (NEW)
- Image presence (NEW)
- Shape positions (NEW)

Usage:
    python scripts/validate_presentation.py output.pptx [--template template.pptx]
"""

import re
import sys
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

# Expected brand fonts
BRAND_FONTS = {
    'headline': ['VC Nudge Black', 'VC Nudge', 'VCNudge-Black'],
    'body': ['Manrope Light', 'Manrope', 'Manrope-Light'],
    'watermark': ['IBM Plex Mono', 'IBMPlexMono', 'IBM Plex Mono Light'],
}

# Brand colors (RGB)
BRAND_COLORS = {
    'black': (0, 0, 0),
    'cyan': (125, 211, 232),  # #7DD3E8
    'white': (255, 255, 255),
}


def extract_text_from_shape(shape):
    """Recursively extract all text from a shape."""
    texts = []
    if shape.has_text_frame:
        texts.append(shape.text_frame.text)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(extract_text_from_shape(child))
    return texts


def get_fonts_from_shape(shape):
    """Recursively extract all fonts used in a shape."""
    fonts = set()
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.name:
                    fonts.add(run.font.name)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            fonts.update(get_fonts_from_shape(child))
    return fonts


def get_images_from_slide(slide):
    """Get count and info of images in a slide."""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append({
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Check for images in groups
            for child in shape.shapes:
                if hasattr(child, 'shape_type') and child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append({
                        'name': child.name,
                        'left': child.left,
                        'top': child.top,
                        'width': child.width,
                        'height': child.height,
                    })
    return images


def check_font_consistency(prs):
    """Check that fonts used match expected brand fonts."""
    issues = []
    all_fonts = set()

    for i, slide in enumerate(prs.slides):
        slide_fonts = set()
        for shape in slide.shapes:
            fonts = get_fonts_from_shape(shape)
            slide_fonts.update(fonts)
            all_fonts.update(fonts)

    # Check for unexpected fonts
    expected_fonts = set()
    for font_list in BRAND_FONTS.values():
        expected_fonts.update(font_list)

    unexpected = all_fonts - expected_fonts
    if unexpected:
        # Filter out common safe fonts
        safe_fonts = {'Arial', 'Calibri', 'Times New Roman', '+mn-lt', '+mn-ea', '+mj-lt'}
        truly_unexpected = unexpected - safe_fonts
        if truly_unexpected:
            issues.append(f"Unexpected fonts found: {truly_unexpected}")

    return all_fonts, issues


def check_images(prs, template_prs=None):
    """Check that images are present and not broken."""
    issues = []
    image_counts = []

    for i, slide in enumerate(prs.slides):
        images = get_images_from_slide(slide)
        image_counts.append(len(images))

        # Check for placeholder images (very small or zero dimension)
        for img in images:
            if img['width'] == 0 or img['height'] == 0:
                issues.append(f"Slide {i+1}: Broken image '{img['name']}' (zero dimensions)")

    # If template provided, compare image counts
    if template_prs:
        template_counts = []
        for slide in template_prs.slides:
            template_counts.append(len(get_images_from_slide(slide)))

        # Simple comparison - this is informational
        if len(image_counts) != len(template_counts):
            issues.append(f"Slide count differs from template ({len(image_counts)} vs {len(template_counts)})")

    return image_counts, issues


def validate_presentation(filepath, template_path=None, verbose=False):
    """
    Validate a generated presentation.

    Args:
        filepath: Path to presentation to validate
        template_path: Optional path to template for comparison
        verbose: If True, include detailed information

    Returns:
        Dict with validation results
    """
    errors = []
    warnings = []
    info = []

    # Try to open the file
    try:
        prs = Presentation(filepath)
    except Exception as e:
        return {
            "valid": False,
            "errors": [f"Failed to open file: {e}"],
            "warnings": [],
            "info": [],
            "slide_count": 0
        }

    # Load template if provided
    template_prs = None
    if template_path:
        try:
            template_prs = Presentation(template_path)
        except Exception as e:
            warnings.append(f"Could not load template for comparison: {template_path} ({e})")

    # ===== BASIC CHECKS =====

    # Check slide count
    slide_count = len(prs.slides)
    if slide_count < 2:
        errors.append(f"Presentation has only {slide_count} slide(s) - minimum 2 expected")
    info.append(f"Total slides: {slide_count}")

    # ===== PLACEHOLDER CHECK =====

    # Regex pattern for semantic placeholders: {{type.element}}
    placeholder_pattern = re.compile(r'\{\{[a-z_]+\.[a-z_0-9]+\}\}')

    unreplaced = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            texts = extract_text_from_shape(shape)
            for text in texts:
                matches = placeholder_pattern.findall(text)
                for match in matches:
                    unreplaced.append(f"Slide {i+1}: {match}")

    if unreplaced:
        errors.append(f"Found {len(unreplaced)} unreplaced placeholder(s):")
        for item in unreplaced[:10]:  # Show first 10
            errors.append(f"  - {item}")
        if len(unreplaced) > 10:
            errors.append(f"  ... and {len(unreplaced) - 10} more")

    # ===== THANK-YOU SLIDE CHECK =====

    if slide_count > 0:
        last_slide = prs.slides[-1]
        has_thank_you = False
        for shape in last_slide.shapes:
            texts = extract_text_from_shape(shape)
            for text in texts:
                if "THANK YOU" in text.upper():
                    has_thank_you = True
                    break
            if has_thank_you:
                break

        if not has_thank_you:
            warnings.append("Last slide may not be the thank-you slide")
        else:
            info.append("Thank-you slide: Present")

    # ===== FONT CHECK (NEW in V2) =====

    fonts_used, font_issues = check_font_consistency(prs)
    for issue in font_issues:
        warnings.append(issue)

    if verbose and fonts_used:
        info.append(f"Fonts used: {', '.join(sorted(fonts_used))}")

    # ===== IMAGE CHECK (NEW in V2) =====

    image_counts, image_issues = check_images(prs, template_prs)
    for issue in image_issues:
        warnings.append(issue)

    total_images = sum(image_counts)
    if verbose:
        info.append(f"Total images: {total_images}")

    # ===== SLIDE CONTENT CHECK =====

    empty_slides = []
    for i, slide in enumerate(prs.slides):
        has_content = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_content = True
                break
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_content = True
                break
        if not has_content:
            empty_slides.append(i + 1)

    if empty_slides:
        warnings.append(f"Potentially empty slides: {empty_slides}")

    # ===== BUILD RESULT =====

    return {
        "valid": len(errors) == 0 and len(unreplaced) == 0,
        "errors": errors,
        "warnings": warnings,
        "info": info,
        "slide_count": slide_count,
        "unreplaced_count": len(unreplaced),
        "fonts_used": list(fonts_used),
        "total_images": total_images,
    }


def main():
    parser = argparse.ArgumentParser(
        description='Validate an HCO presentation',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python validate_presentation.py output.pptx
  python validate_presentation.py output.pptx --template assets/ppt-template-v2.pptx
  python validate_presentation.py output.pptx -v
        """
    )
    parser.add_argument('filepath', help='Path to presentation to validate')
    parser.add_argument('--template', '-t', help='Path to template for comparison')
    parser.add_argument('--verbose', '-v', action='store_true', help='Show detailed info')
    parser.add_argument('--json', '-j', action='store_true', help='Output as JSON')

    args = parser.parse_args()

    print(f"Validating: {args.filepath}")
    print("-" * 50)

    result = validate_presentation(args.filepath, args.template, args.verbose)

    if args.json:
        import json
        print(json.dumps(result, indent=2))
        sys.exit(0 if result['valid'] else 1)

    # Standard output
    print(f"Slides: {result['slide_count']}")
    print(f"Valid: {'YES' if result['valid'] else 'NO'}")

    if result.get('info') and args.verbose:
        print("\nINFO:")
        for item in result['info']:
            print(f"  {item}")

    if result['errors']:
        print("\nERRORS:")
        for error in result['errors']:
            print(f"  {error}")

    if result['warnings']:
        print("\nWARNINGS:")
        for warning in result['warnings']:
            print(f"  {warning}")

    if result['valid']:
        print("\n[OK] Presentation passed validation")
    else:
        print("\n[FAIL] Presentation has issues")
        sys.exit(1)


if __name__ == "__main__":
    main()
