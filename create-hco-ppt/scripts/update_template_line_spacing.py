#!/usr/bin/env python3
"""
Apply 0.7 line spacing to all VC Nudge Black headlines in the template.

This script updates the PPT template to use tighter line spacing (0.7× multiple)
for all VC Nudge Black headline text. This is a one-time update script.

Usage:
    cd skills/create-hco-ppt-v2
    python scripts/update_template_line_spacing.py
"""

import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Line spacing value (0.7 = 70% of single line spacing)
VC_NUDGE_LINE_SPACING = 0.7

# All slides that contain VC Nudge headlines
VC_NUDGE_SLIDE_INDICES = [0, 2, 4, 5, 6, 7, 8, 9, 10, 11, 12, 14, 15, 16, 18, 19, 20, 21, 23]


def is_vc_nudge_font(para):
    """
    Check if paragraph uses VC Nudge font.

    Args:
        para: A pptx paragraph object

    Returns:
        True if any run in the paragraph uses VC Nudge font
    """
    for run in para.runs:
        font_name = run.font.name
        if font_name and 'VC Nudge' in font_name:
            return True
    return False


def apply_line_spacing_to_shape(shape, count=0):
    """
    Recursively apply 0.7 line spacing to VC Nudge text in a shape.

    Args:
        shape: A pptx shape object
        count: Running count of updated paragraphs

    Returns:
        Updated count of paragraphs modified
    """
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if is_vc_nudge_font(para):
                para.line_spacing = VC_NUDGE_LINE_SPACING
                count += 1

    # Handle grouped shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            count = apply_line_spacing_to_shape(child, count)

    return count


def main():
    """Update template with 0.7 line spacing for VC Nudge headlines."""
    # Determine paths
    script_dir = os.path.dirname(os.path.abspath(__file__))
    skill_dir = os.path.dirname(script_dir)
    template_path = os.path.join(skill_dir, 'assets', 'ppt-template-v2.pptx')

    if not os.path.exists(template_path):
        print(f"Error: Template not found at {template_path}")
        sys.exit(1)

    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)

    total_updated = 0
    slides_updated = 0

    for idx in VC_NUDGE_SLIDE_INDICES:
        if idx >= len(prs.slides):
            print(f"Warning: Slide index {idx} exceeds template size, skipping")
            continue

        slide = prs.slides[idx]
        slide_count = 0

        for shape in slide.shapes:
            slide_count += apply_line_spacing_to_shape(shape)

        if slide_count > 0:
            print(f"  Slide {idx + 1} (index {idx}): {slide_count} paragraph(s) updated")
            total_updated += slide_count
            slides_updated += 1

    # Save the updated template
    prs.save(template_path)

    print(f"\nComplete: Updated {total_updated} VC Nudge paragraphs across {slides_updated} slides")
    print(f"Template saved: {template_path}")


if __name__ == "__main__":
    main()
